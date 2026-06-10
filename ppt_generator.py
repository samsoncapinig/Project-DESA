from pptx import Presentation
from datetime import datetime
import re


# ===== NORMALIZE FUNCTION =====
def _normalize(text):
    return re.sub(r'[^a-z0-9 ]', '', text.lower().replace("_", " "))


# ===== MAIN FUNCTION =====
def generate_ppt_report(data, daily_rating_data, end_rating_data):

    prs = Presentation("desa_template.pptx")

    # ===== SLIDE 1 =====
    if len(prs.slides) > 0:
        slide1 = prs.slides[0]
        _replace_text_in_slide(slide1, {
            "{{title}}": data.get("training_title", ""),
            "{{date}}": data.get("date", ""),
            "{{venue}}": data.get("venue", ""),
            "{{program_owner}}": data.get("program_owner", ""),
        })

    # ===== SLIDE 2 =====
    if len(prs.slides) > 1:
        slide2 = prs.slides[1]

        program_management = _find_matching_value(daily_rating_data, ["program", "management"])
        accommodation = _find_matching_value(daily_rating_data, ["accommodation"])
        training_venue = _find_matching_value(daily_rating_data, ["training", "venue",])
        food = _find_matching_value(daily_rating_data, ["food"])
        administrative_arrangements = _find_matching_value(daily_rating_data, ["administrative"])

def _is_category_rating(key):
    key = key.lower()
    return any(x in key for x in [
        "program", "accommodation", "venue", "food", "administrative"
    ])

category_ratings = [
    float(v) for k, v in average_rating.items()
    if _is_average_rating(k) and isinstance(v, (int, float))
]

avg_sessions = sum(category_ratings) / len(category_ratings) if category_ratings else 0


        _replace_text_in_slide(slide2, {
            "{{program_management}}": _format_value(program_management),
            "{{accommodation}}": _format_value(accommodation),
            "{{training_venue}}": _format_value(training_venue),
            "{{food}}": _format_value(food),
            "{{administrative_arrangements}}": _format_value(administrative_arrangements),
            "{{overall_daily_average}}": _format_value(data.get("daily_general_average", avg_sessions)),
        })

    # ===== SLIDE 3 =====
# ===== SLIDE 3 =====
if len(prs.slides) > 2:
    slide3 = prs.slides[2]

    # ✅ COMPUTE END PROGRAM AVERAGE (Fix 4 here)
    end_values = [
        float(v) for k, v in end_rating_data.items()
        if isinstance(v, (int, float))
    ]

    end_avg = sum(end_values) / len(end_values) if end_values else 0

    replacements = {
        "{{program_management1}}": _format_value(_find_matching_value(end_rating_data, ["program", "management"])),
        "{{attainment_of_objectives}}": _format_value(_find_matching_value(end_rating_data, ["attainment"])),
        "{{delivery_of_content}}": _format_value(_find_matching_value(end_rating_data, ["delivery", "content"])),
        "{{provision_of_support_materials}}": _format_value(_find_matching_value(end_rating_data, ["support", "materials"])),
        "{{program_management_team}}": _format_value(_find_matching_value(end_rating_data, ["program", "team"])),
        "{{training_venue1}}": _format_value(_find_matching_value(end_rating_data, ["training", "venue"])),
        "{{food1}}": _format_value(_find_matching_value(end_rating_data, ["food"])),
        "{{accommodation1}}": _format_value(_find_matching_value(end_rating_data, ["accommodation"])),

        # ✅ USE COMPUTED VALUE (not data.get anymore)
        "{{end_of_program_average}}": _format_value(end_avg),
    }

    # ===== SESSION SLIDES =====
    session_map = [
        ("day1", 3),
        ("day2", 4),
        ("day3", 5),
        ("day4", 6),
        ("day5", 7),
    ]

    for day, slide_index in session_map:
        if len(prs.slides) > slide_index:
            slide = prs.slides[slide_index]

            replacements = {}
            for i in range(1, 6):
                key = f"{day}_lm{i}"
                value = _find_matching_value(daily_rating_data, [day, f"lm{i}"])
                replacements[f"{{{{{key}}}}}"] = _format_value(value)

            _replace_text_in_slide(slide, replacements)

    # ===== SLIDE 9 =====
    if len(prs.slides) > 8:
        _replace_text_in_slide(prs.slides[8], {
            "{{analysis}}": data.get("analysis", "")
        })

    # ===== SLIDE 10 =====
    if len(prs.slides) > 9:
        _replace_text_in_slide(prs.slides[9], {
            "{{recommendation}}": data.get("recommendation", "")
        })

    # ===== SLIDE 11 =====
    if len(prs.slides) > 10:
        _replace_text_in_slide(prs.slides[10], {
            "{{daily_general_average}}": _format_value(data.get("daily_general_average", 0)),
            "{{end_of_program_average}}": _format_value(data.get("end_of_program_average", 0)),
            "{{overall_results}}": _format_value(data.get("overall_results", 0)),
        })

    # ===== SLIDE 12 =====
    if len(prs.slides) > 11:
        _replace_text_in_slide(prs.slides[11], {
            "{{overall_results}}": _format_value(data.get("overall_results", 0)),
        })

    filename = f"DESA_Report_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
    prs.save(filename)

    return filename


# ===== HELPERS =====

def _find_matching_value(data_dict, keywords):
    if not data_dict:
        return "N/A"

    keywords = [_normalize(k) for k in keywords]

    best_match = None

    for key, value in data_dict.items():
        key_clean = _normalize(key)

        match_count = sum(1 for k in keywords if k in key_clean)

        if match_count > 0:
            # prioritize strongest match
            if best_match is None or match_count > best_match[1]:
                best_match = (value, match_count)

    return best_match[0] if best_match else "N/A"

def _is_session_rating(key):
    key = key.lower().replace("_", " ")
    return "day" in key and "lm" in key


def _replace_text_in_slide(slide, replacements):
    for shape in slide.shapes:

        # TEXT
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                text = "".join(run.text for run in paragraph.runs)

                for placeholder, value in replacements.items():
                    text = text.replace(placeholder, str(value))

                paragraph.text = text

        # ✅ TABLES (FIXED VERSION)
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    text = cell.text

                    for placeholder, value in replacements.items():
                        text = text.replace(placeholder, str(value))

                    cell.text = text


def _format_value(value):
    if isinstance(value, (int, float)):
        return f"{value:.2f}"
    return str(value) if value else "N/A"
